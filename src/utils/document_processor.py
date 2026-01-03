# ==============================================================================
# 文档处理器模块
# ==============================================================================
# 提供多种文档格式的解析和转换功能
#
# 支持格式
# - PDF: 文本提取和图片提取
# - Word (.docx/.doc): 文本和图片提取
# - PowerPoint (.pptx): 幻灯片内容提取
# - Excel (.xlsx/.xls): 表格数据读取
# - 图片 (png/jpg/gif等): Base64编码和描述
# - 纯文本 (txt/md): 直接读取
#
# 基于 GAIA Benchmark 需求设计
# ==============================================================================

import base64
import csv
import os
import re
import shutil
import time
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from src.core.logger import get_logger
from src.utils.encoding import decode_output_bytes

logger = get_logger("utils.document_processor")

# 尝试导入可选依赖
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.warning("Pillow 未安装，图片处理功能受限 / Pillow not installed, image processing limited")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# pywin32 用于Windows下读取老式.doc文件
try:
    import win32com.client
    import pythoncom
    PYWIN32_AVAILABLE = True
except ImportError:
    PYWIN32_AVAILABLE = False

# olefile 用于跨平台读取老式.doc文件
try:
    import olefile
    OLEFILE_AVAILABLE = True
except ImportError:
    OLEFILE_AVAILABLE = False

# subprocess 用于调用外部工具
import subprocess
import platform


# ==============================================================================
# 辅助函数
# ==============================================================================

def encode_image_to_base64(image_path: str) -> str:
    """
    将图片文件编码为Base64字符串 / Encode image file to Base64 string
    
    Args:
        image_path: 图片文件路径 / Image file path
        
    Returns:
        str: Base64编码的字符串 / Base64 encoded string
    """
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')


def get_image_size(image_path: str) -> Tuple[int, int]:
    """
    获取图片尺寸 / Get image dimensions
    
    Args:
        image_path: 图片文件路径 / Image file path
        
    Returns:
        Tuple[int, int]: (宽度, 高度) / (width, height)
        
    Raises:
        FileNotFoundError: 文件不存在 / File not found
        ValueError: 无法读取图片 / Cannot read image
    """
    if not PIL_AVAILABLE:
        raise ImportError("Pillow 未安装 / Pillow is not installed")
    
    try:
        with Image.open(image_path) as img:
            return img.size
    except FileNotFoundError:
        raise FileNotFoundError(f"图片文件未找到 / Image file not found: {image_path}")
    except Exception as e:
        raise ValueError(f"无法读取图片 / Cannot read image {image_path}: {str(e)}")


def check_dependencies(file_ext: str) -> Optional[str]:
    """
    检查处理特定文件类型所需的依赖 / Check dependencies for specific file type
    
    Args:
        file_ext: 文件扩展名（不含点）/ File extension (without dot)
        
    Returns:
        Optional[str]: 缺失的依赖名称，None表示依赖已满足 / Missing dependency name, None if satisfied
    """
    ext = file_ext.lower().lstrip('.')
    
    if ext == 'pdf':
        if not PDFPLUMBER_AVAILABLE:
            return "pdfplumber"
    elif ext == 'docx':
        if not DOCX_AVAILABLE:
            return "python-docx"
    elif ext == 'doc':
        # .doc文件需要pywin32(Windows)或olefile(跨平台)
        if not PYWIN32_AVAILABLE and not OLEFILE_AVAILABLE:
            if platform.system() == 'Windows':
                return "pywin32"
            else:
                return "olefile"
    elif ext == 'pptx':
        if not PPTX_AVAILABLE:
            return "python-pptx"
    elif ext in ['ppt']:
        # 老式.ppt文件需要pywin32
        if not PYWIN32_AVAILABLE:
            if platform.system() == 'Windows':
                return "pywin32"
            else:
                return "libreoffice (soffice command)"
    elif ext in ['xlsx', 'xls']:
        if not PANDAS_AVAILABLE:
            return "pandas"
    elif ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp']:
        if not PIL_AVAILABLE:
            return "Pillow"
    
    return None


# ==============================================================================
# 文档转换器类
# ==============================================================================

class DocumentConverter:
    """
    文档转换器 / Document Converter
    
    将各种文档格式转换为统一的文本/Markdown格式，便于Agent处理
    Converts various document formats to unified text/Markdown format for Agent processing
    
    使用方法 / Usage:
        converter = DocumentConverter(task_id="task_001", output_dir="./temp")
        result = converter.convert("document.pdf")
    """
    
    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        'document': ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.md'],
        'image': ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.svg'],
        'spreadsheet': ['.xlsx', '.xls', '.csv'],
    }
    
    def __init__(
        self,
        task_id: str,
        output_dir: Optional[str] = None,
        image_description_func: Optional[callable] = None
    ):
        """
        初始化文档转换器 / Initialize document converter
        
        Args:
            task_id: 任务ID，用于组织输出目录 / Task ID for organizing output directory
            output_dir: 输出目录，默认为 ./workspace/temp/{task_id}/input
                       Output directory, default is ./workspace/temp/{task_id}/input
            image_description_func: 图片描述函数，用于生成图片的文字描述
                                   Image description function for generating text descriptions
        """
        self.task_id = task_id
        
        # 设置输出目录
        if output_dir:
            self.output_dir = Path(output_dir) / task_id / "input"
        else:
            self.output_dir = Path("./workspace/temp") / task_id / "input"
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # 图片CSV路径
        self.image_csv_path = self.output_dir / "images.csv"
        
        # 图片计数器
        self.image_counter = 1
        
        # 图片描述函数
        self.image_description_func = image_description_func
        
        # 初始化图片CSV
        self._init_image_csv()
        
        logger.info(f"文档转换器初始化 / Document converter initialized: {self.output_dir}")
    
    def _init_image_csv(self) -> None:
        """初始化图片信息CSV文件 / Initialize image info CSV file"""
        # 如果文件已存在，删除重建
        if self.image_csv_path.exists():
            self.image_csv_path.unlink()
        
        with open(self.image_csv_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerow(["name", "source_path", "base64", "description", "size"])
    
    def _save_image_info(
        self,
        name: str,
        source: str,
        base64_str: str,
        description: str,
        size: Union[str, Tuple[int, int]]
    ) -> None:
        """
        保存图片信息到CSV / Save image info to CSV
        
        Args:
            name: 图片名称 / Image name
            source: 来源路径 / Source path
            base64_str: Base64编码 / Base64 encoding
            description: 图片描述 / Image description
            size: 尺寸（宽x高或字符串）/ Size (WxH or string)
        """
        if isinstance(size, tuple):
            size_str = f"{size[0]}x{size[1]}"
        else:
            size_str = str(size)
        
        with open(self.image_csv_path, 'a', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerow([name, source, base64_str, description, size_str])
    
    def _image_to_base64(self, image_data: bytes) -> Tuple[str, str]:
        """
        将图片数据转换为Base64 / Convert image data to Base64
        
        Args:
            image_data: 图片二进制数据 / Image binary data
            
        Returns:
            Tuple[str, str]: (Base64字符串, 尺寸字符串) / (Base64 string, size string)
        """
        if not PIL_AVAILABLE:
            # 没有PIL时，直接编码原始数据
            return base64.b64encode(image_data).decode('utf-8'), "unknown"
        
        try:
            img = Image.open(BytesIO(image_data))
            buffered = BytesIO()
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode('utf-8'), f"{img.size[0]}x{img.size[1]}"
        except Exception as e:
            # 处理损坏或无法识别的图片
            logger.warning(f"无法处理图片，使用占位符 / Cannot process image, using placeholder: {e}")
            # 创建占位图片
            placeholder = Image.new('RGB', (100, 100), color='gray')
            buffered = BytesIO()
            placeholder.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode('utf-8'), "100x100"
    
    # ==========================================================================
    # 各类型文件转换方法
    # ==========================================================================
    
    def convert_pdf(self, file_path: str) -> str:
        """
        转换PDF文件 / Convert PDF file
        
        提取PDF中的文本和图片
        Extracts text and images from PDF
        
        Args:
            file_path: PDF文件路径 / PDF file path
            
        Returns:
            str: Markdown格式的内容 / Markdown formatted content
        """
        if not PDFPLUMBER_AVAILABLE:
            raise ImportError("pdfplumber 未安装，请运行: pip install pdfplumber")
        
        md_content = []
        file_name = Path(file_path).stem
        
        # 使用pdfplumber提取文字
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    md_content.append(f"--- Page {page_num} ---\n{text}")
                
                # 提取表格
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables, 1):
                    if table:
                        md_content.append(f"\n[表格 {page_num}-{table_idx}]\n")
                        # 简单的表格转文本
                        for row in table:
                            md_content.append(" | ".join(str(cell or '') for cell in row))
        
        # 使用PyMuPDF提取图片（如果可用）
        if FITZ_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    image_list = page.get_images(full=True)
                    
                    for img_index, img in enumerate(image_list, 1):
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            image_ext = base_image["ext"]
                            
                            base64_str, size = self._image_to_base64(image_bytes)
                            img_name = f"{file_name}_page{page_num + 1}_img{img_index}.{image_ext}"
                            self._save_image_info(
                                img_name,
                                file_path,
                                base64_str,
                                f"从PDF第{page_num + 1}页提取 / Extracted from PDF page {page_num + 1}",
                                size
                            )
                            self.image_counter += 1
                        except Exception as e:
                            logger.warning(f"PDF图片提取失败 / Failed to extract PDF image: {e}")
                doc.close()
            except Exception as e:
                logger.warning(f"PyMuPDF处理失败 / PyMuPDF processing failed: {e}")
        
        return "\n\n".join(md_content)
    
    def convert_word(self, file_path: str) -> str:
        """
        转换Word文档(.docx) / Convert Word document (.docx)
        
        Args:
            file_path: Word文件路径 / Word file path
            
        Returns:
            str: Markdown格式的内容 / Markdown formatted content
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx 未安装，请运行: pip install python-docx")
        
        doc = docx.Document(file_path)
        file_name = Path(file_path).stem
        md_content = []
        
        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                md_content.append(para.text)
        
        # 提取表格
        for table_idx, table in enumerate(doc.tables, 1):
            md_content.append(f"\n[表格 {table_idx}]")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                md_content.append(row_text)
        
        # 提取图片
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_data = rel.target_part.blob
                    base64_str, size = self._image_to_base64(img_data)
                    img_name = f"{file_name}_word_image_{self.image_counter}"
                    self._save_image_info(
                        img_name,
                        file_path,
                        base64_str,
                        "从Word文档提取 / Extracted from Word document",
                        size
                    )
                    self.image_counter += 1
                except Exception as e:
                    logger.warning(f"Word图片提取失败 / Failed to extract Word image: {e}")
        
        return "\n\n".join(md_content)
    
    def convert_doc_legacy(self, file_path: str) -> str:
        """
        转换老式Word文档(.doc) / Convert legacy Word document (.doc)
        
        支持多种方法按优先级尝试 / Supports multiple methods with priority:
        1. pywin32 (仅Windows，最可靠) / pywin32 (Windows only, most reliable)
        2. antiword (跨平台命令行工具) / antiword (cross-platform CLI tool)
        3. olefile (跨平台，基本文本提取) / olefile (cross-platform, basic text extraction)
        
        Args:
            file_path: .doc文件路径 / .doc file path
            
        Returns:
            str: 提取的文本内容 / Extracted text content
        """
        file_name = Path(file_path).stem
        md_content = []
        
        # 方法1: 使用 pywin32 (Windows COM接口)
        if PYWIN32_AVAILABLE and platform.system() == 'Windows':
            try:
                pythoncom.CoInitialize()
                try:
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    
                    # 打开文档
                    doc = word.Documents.Open(os.path.abspath(file_path))
                    
                    # 提取文本
                    text = doc.Content.Text
                    md_content.append(text)
                    
                    # 提取表格
                    for table_idx, table in enumerate(doc.Tables, 1):
                        md_content.append(f"\n[表格 {table_idx}]")
                        try:
                            for row in table.Rows:
                                row_texts = []
                                for cell in row.Cells:
                                    row_texts.append(cell.Range.Text.strip().replace('\r\x07', ''))
                                md_content.append(" | ".join(row_texts))
                        except Exception as e:
                            logger.warning(f"表格提取部分失败 / Table extraction partially failed: {e}")
                    
                    doc.Close(False)
                    word.Quit()
                    
                    logger.info(f"使用pywin32成功转换.doc文件 / Successfully converted .doc file using pywin32: {file_path}")
                    return "\n\n".join(md_content)
                    
                finally:
                    pythoncom.CoUninitialize()
                    
            except Exception as e:
                logger.warning(f"pywin32转换失败，尝试其他方法 / pywin32 conversion failed, trying other methods: {e}")
        
        # 方法2: 使用 antiword 命令行工具
        try:
            antiword_cmd = "antiword" if platform.system() != 'Windows' else "antiword.exe"
            result = subprocess.run(
                [antiword_cmd, file_path],
                capture_output=True,
                timeout=30,
                text=False
            )
            stdout_text = decode_output_bytes(result.stdout) if result.stdout else ""
            if result.returncode == 0 and stdout_text.strip():
                logger.info(f"使用antiword成功转换.doc文件 / Successfully converted .doc file using antiword: {file_path}")
                return stdout_text
        except FileNotFoundError:
            logger.debug("antiword未安装 / antiword not installed")
        except subprocess.TimeoutExpired:
            logger.warning("antiword执行超时 / antiword execution timeout")
        except Exception as e:
            logger.warning(f"antiword转换失败 / antiword conversion failed: {e}")
        
        # 方法3: 使用 catdoc 命令行工具 (常见于Linux)
        try:
            result = subprocess.run(
                ["catdoc", file_path],
                capture_output=True,
                timeout=30,
                text=False
            )
            stdout_text = decode_output_bytes(result.stdout) if result.stdout else ""
            if result.returncode == 0 and stdout_text.strip():
                logger.info(f"使用catdoc成功转换.doc文件 / Successfully converted .doc file using catdoc: {file_path}")
                return stdout_text
        except FileNotFoundError:
            logger.debug("catdoc未安装 / catdoc not installed")
        except Exception as e:
            logger.warning(f"catdoc转换失败 / catdoc conversion failed: {e}")
        
        # 方法4: 使用 olefile 进行基本文本提取
        if OLEFILE_AVAILABLE:
            try:
                ole = olefile.OleFileIO(file_path)
                
                # 尝试读取 WordDocument 流
                if ole.exists('WordDocument'):
                    # 读取所有文本流
                    text_parts = []
                    
                    # 尝试从多个可能的位置提取文本
                    for stream_name in ole.listdir():
                        stream_path = '/'.join(stream_name)
                        try:
                            if any(keyword in stream_path.lower() for keyword in ['word', 'text', 'content', 'data']):
                                data = ole.openstream(stream_name).read()
                                # 尝试提取可读文本
                                text = self._extract_text_from_binary(data)
                                if text.strip():
                                    text_parts.append(text)
                        except Exception:
                            continue
                    
                    ole.close()
                    
                    if text_parts:
                        logger.info(f"使用olefile成功转换.doc文件 / Successfully converted .doc file using olefile: {file_path}")
                        return "\n\n".join(text_parts)
                
                ole.close()
                
            except Exception as e:
                logger.warning(f"olefile转换失败 / olefile conversion failed: {e}")
        
        # 方法5: 最后尝试直接二进制文本提取
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
            text = self._extract_text_from_binary(raw_data)
            if text.strip():
                logger.info(f"使用二进制提取成功转换.doc文件 / Successfully converted .doc file using binary extraction: {file_path}")
                return text
        except Exception as e:
            logger.warning(f"二进制文本提取失败 / Binary text extraction failed: {e}")
        
        raise ValueError(
            f"无法读取.doc文件 / Cannot read .doc file: {file_path}\n"
            f"请安装以下工具之一 / Please install one of the following:\n"
            f"- Windows: pip install pywin32 (需要安装Microsoft Word / requires MS Word)\n"
            f"- Linux/Mac: sudo apt-get install antiword 或 catdoc\n"
            f"- 跨平台 / Cross-platform: pip install olefile"
        )
    
    def _extract_text_from_binary(self, data: bytes) -> str:
        """
        从二进制数据中提取可读文本 / Extract readable text from binary data
        
        Args:
            data: 二进制数据 / Binary data
            
        Returns:
            str: 提取的文本 / Extracted text
        """
        # 尝试多种编码
        encodings = ['utf-16-le', 'utf-8', 'cp1252', 'gbk', 'latin-1']
        
        for encoding in encodings:
            try:
                text = data.decode(encoding, errors='ignore')
                # 过滤不可打印字符，保留中文和常见字符
                import re
                # 保留可打印ASCII、中文、常见标点
                cleaned = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\n\r\t]', ' ', text)
                # 压缩多余空格
                cleaned = re.sub(r' {3,}', '  ', cleaned)
                cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
                
                # 检查是否有足够的可读内容
                if len(cleaned.strip()) > 50:
                    return cleaned.strip()
            except Exception:
                continue
        
        return ""
    
    def convert_ppt_legacy(self, file_path: str) -> str:
        """
        转换老式PowerPoint文档(.ppt) / Convert legacy PowerPoint document (.ppt)
        
        Args:
            file_path: .ppt文件路径 / .ppt file path
            
        Returns:
            str: 提取的文本内容 / Extracted text content
        """
        file_name = Path(file_path).stem
        md_content = []
        
        # 方法1: 使用 pywin32 (Windows COM接口)
        if PYWIN32_AVAILABLE and platform.system() == 'Windows':
            try:
                pythoncom.CoInitialize()
                try:
                    ppt = win32com.client.Dispatch("PowerPoint.Application")
                    ppt.Visible = False
                    
                    # 打开演示文稿
                    presentation = ppt.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                    
                    # 遍历幻灯片
                    for slide_idx, slide in enumerate(presentation.Slides, 1):
                        md_content.append(f"\n--- Slide {slide_idx} ---\n")
                        
                        # 提取形状中的文本
                        for shape in slide.Shapes:
                            if shape.HasTextFrame:
                                if shape.TextFrame.HasText:
                                    text = shape.TextFrame.TextRange.Text
                                    if text.strip():
                                        md_content.append(text)
                    
                    presentation.Close()
                    ppt.Quit()
                    
                    logger.info(f"使用pywin32成功转换.ppt文件 / Successfully converted .ppt file using pywin32: {file_path}")
                    return "\n\n".join(md_content)
                    
                finally:
                    pythoncom.CoUninitialize()
                    
            except Exception as e:
                logger.warning(f"pywin32转换.ppt失败，尝试其他方法 / pywin32 .ppt conversion failed, trying other methods: {e}")
        
        # 方法2: 使用 olefile 进行基本文本提取
        if OLEFILE_AVAILABLE:
            try:
                ole = olefile.OleFileIO(file_path)
                
                text_parts = []
                for stream_name in ole.listdir():
                    stream_path = '/'.join(stream_name)
                    try:
                        if any(keyword in stream_path.lower() for keyword in ['powerpoint', 'text', 'current']):
                            data = ole.openstream(stream_name).read()
                            text = self._extract_text_from_binary(data)
                            if text.strip():
                                text_parts.append(text)
                    except Exception:
                        continue
                
                ole.close()
                
                if text_parts:
                    logger.info(f"使用olefile成功转换.ppt文件 / Successfully converted .ppt file using olefile: {file_path}")
                    return "\n\n".join(text_parts)
                    
            except Exception as e:
                logger.warning(f"olefile转换.ppt失败 / olefile .ppt conversion failed: {e}")
        
        raise ValueError(
            f"无法读取.ppt文件 / Cannot read .ppt file: {file_path}\n"
            f"请安装以下工具之一 / Please install one of the following:\n"
            f"- Windows: pip install pywin32 (需要安装Microsoft PowerPoint / requires MS PowerPoint)\n"
            f"- 跨平台 / Cross-platform: pip install olefile"
        )
    
    def convert_ppt(self, file_path: str) -> str:
        """
        转换PowerPoint文件 / Convert PowerPoint file
        
        Args:
            file_path: PPT文件路径 / PPT file path
            
        Returns:
            str: Markdown格式的内容 / Markdown formatted content
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")
        
        prs = pptx.Presentation(file_path)
        file_name = Path(file_path).stem
        md_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            md_content.append(f"\n--- Slide {slide_num} ---\n")
            
            for shape in slide.shapes:
                # 提取文本
                if hasattr(shape, "text") and shape.text:
                    md_content.append(shape.text)
                
                # 提取图片
                if shape.shape_type == 13:  # Picture shape
                    try:
                        img_data = shape.image.blob
                        base64_str, size = self._image_to_base64(img_data)
                        img_name = f"{file_name}_slide{slide_num}_img_{self.image_counter}"
                        self._save_image_info(
                            img_name,
                            file_path,
                            base64_str,
                            f"从幻灯片第{slide_num}页提取 / Extracted from slide {slide_num}",
                            size
                        )
                        self.image_counter += 1
                    except Exception as e:
                        logger.warning(f"PPT图片提取失败 / Failed to extract PPT image: {e}")
        
        return "\n\n".join(md_content)
    
    def convert_excel(self, file_path: str) -> Tuple[str, str]:
        """
        转换Excel文件 / Convert Excel file
        
        Args:
            file_path: Excel文件路径 / Excel file path
            
        Returns:
            Tuple[str, str]: (内容预览, 复制后的文件路径) / (content preview, copied file path)
        """
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas 未安装，请运行: pip install pandas openpyxl")
        
        # 复制文件到输出目录
        output_path = self.output_dir / Path(file_path).name
        shutil.copy2(file_path, output_path)
        
        # 读取所有工作表
        try:
            xls = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name, nrows=10)
                content_parts.append(f"=== Sheet: {sheet_name} ===")
                content_parts.append(f"行数: {len(df)}, 列数: {len(df.columns)}")
                content_parts.append(df.to_markdown(index=False))
            
            content = "\n\n".join(content_parts)
        except Exception as e:
            # 尝试CSV格式
            try:
                df = pd.read_csv(file_path, nrows=10)
                content = f"CSV文件预览 (前10行):\n{df.to_markdown(index=False)}"
            except Exception as e2:
                content = f"无法读取文件: {e}"
        
        return content, str(output_path)
    
    def convert_txt(self, file_path: str, encoding: str = "utf-8") -> str:
        """
        转换纯文本文件 / Convert plain text file
        
        Args:
            file_path: 文本文件路径 / Text file path
            encoding: 文件编码 / File encoding
            
        Returns:
            str: 文件内容 / File content
        """
        # 尝试多种编码
        encodings = [encoding, 'utf-8', 'gbk', 'gb2312', 'latin-1']
        
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        # 最后尝试二进制读取
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='replace')
    
    def convert_markdown(self, file_path: str) -> str:
        """
        转换Markdown文件 / Convert Markdown file
        
        Args:
            file_path: Markdown文件路径 / Markdown file path
            
        Returns:
            str: 文件内容 / File content
        """
        return self.convert_txt(file_path)
    
    def convert_image(self, file_path: str) -> Tuple[str, str]:
        """
        处理图片文件 / Process image file
        
        复制图片到输出目录，生成Base64编码和描述
        Copy image to output directory, generate Base64 encoding and description
        
        Args:
            file_path: 图片文件路径 / Image file path
            
        Returns:
            Tuple[str, str]: (图片描述, 复制后的文件路径) / (image description, copied file path)
        """
        if not PIL_AVAILABLE:
            raise ImportError("Pillow 未安装，请运行: pip install Pillow")
        
        # 复制图片到输出目录
        output_path = self.output_dir / Path(file_path).name
        shutil.copy2(file_path, output_path)
        
        # 获取图片信息
        file_name = Path(file_path).stem
        ext = Path(file_path).suffix.lstrip('.')
        
        base64_image = encode_image_to_base64(file_path)
        size = get_image_size(file_path)
        
        # 生成图片描述
        if self.image_description_func:
            try:
                description = self.image_description_func(base64_image, ext)
            except Exception as e:
                logger.warning(f"图片描述生成失败 / Image description generation failed: {e}")
                description = f"图片文件 / Image file: {file_name}.{ext}, 尺寸 / Size: {size[0]}x{size[1]}"
        else:
            description = f"图片文件 / Image file: {file_name}.{ext}, 尺寸 / Size: {size[0]}x{size[1]}"
        
        # 保存图片信息
        self._save_image_info(file_name, file_path, base64_image, description, size)
        
        return f"{file_name} 的图片描述: {description}", str(output_path)
    
    # ==========================================================================
    # 主转换方法
    # ==========================================================================
    
    def convert(self, file_path: str) -> Tuple[str, str]:
        """
        转换单个文件 / Convert single file
        
        根据文件扩展名自动选择转换方法
        Automatically selects conversion method based on file extension
        
        Args:
            file_path: 文件路径 / File path
            
        Returns:
            Tuple[str, str]: (转换后的内容, 输出路径) / (converted content, output path)
            
        Raises:
            FileNotFoundError: 文件不存在 / File not found
            ValueError: 不支持的文件类型 / Unsupported file type
        """
        if not file_path or not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在 / File not found: {file_path}")
        
        file_name = Path(file_path).stem
        ext = Path(file_path).suffix.lower()
        
        # 检查依赖
        missing_dep = check_dependencies(ext)
        if missing_dep:
            raise ImportError(f"缺少依赖 / Missing dependency: {missing_dep}. 请运行 / Please run: pip install {missing_dep}")
        
        # 根据扩展名选择转换方法
        if ext == '.pdf':
            content = self.convert_pdf(file_path)
        elif ext == '.docx':
            content = self.convert_word(file_path)
        elif ext == '.doc':
            content = self.convert_doc_legacy(file_path)
        elif ext == '.pptx':
            content = self.convert_ppt(file_path)
        elif ext == '.ppt':
            content = self.convert_ppt_legacy(file_path)
        elif ext == '.txt':
            content = self.convert_txt(file_path)
        elif ext == '.md':
            content = self.convert_markdown(file_path)
        elif ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.svg']:
            return self.convert_image(file_path)
        elif ext in ['.xlsx', '.xls', '.csv']:
            return self.convert_excel(file_path)
        else:
            raise ValueError(f"不支持的文件类型 / Unsupported file type: {ext}")
        
        # 保存转换后的Markdown文件
        output_path = self.output_dir / f"{file_name}.md"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        logger.info(f"文件已转换 / File converted: {file_path} -> {output_path}")
        
        return content, str(output_path)
    
    def process(self, file_list: List[str]) -> Dict[str, Any]:
        """
        批量处理文件列表 / Process file list in batch
        
        Args:
            file_list: 文件路径列表 / List of file paths
            
        Returns:
            Dict[str, Any]: 处理结果 / Processing result
                - file_count: 文档文件数量 / Document file count
                - image_count: 图片文件数量 / Image file count
                - files: 输出文件路径列表 / Output file path list
                - content: 合并的内容 / Merged content
        """
        contents = []
        output_paths = []
        image_count = 0
        
        # 图片扩展名列表
        image_exts = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.svg'}
        
        for file_path in file_list:
            try:
                # 统计图片数量
                ext = Path(file_path).suffix.lower()
                if ext in image_exts:
                    image_count += 1
                
                # 转换文件
                content, output_path = self.convert(file_path)
                output_paths.append(output_path)
                
                # 构建内容摘要，保留结构信息
                excerpt = self._build_context_excerpt(content)
                contents.append(f"文件 {output_path} 的内容:\n{excerpt}")
                
            except Exception as e:
                logger.error(f"文件处理失败 / File processing failed: {file_path}, 错误 / Error: {e}")
                contents.append(f"文件 {file_path} 处理失败: {e}")
        
        # 如果没有图片，删除图片CSV
        if image_count == 0 and self.image_csv_path.exists():
            self.image_csv_path.unlink()
        
        return {
            "file_count": len(output_paths) - image_count,
            "image_count": image_count,
            "files": output_paths,
            "content": "\n\n".join(contents)
        }

    def _build_context_excerpt(self, content: str, max_chars: int = 2000) -> str:
        """
        构建可读的内容摘要，保留首尾结构
        """
        if len(content) <= max_chars:
            return content
        head_len = max_chars // 2
        tail_len = max_chars - head_len
        return content[:head_len] + "\n[... 内容过长已截断 ...]\n" + content[-tail_len:]
    
    def cleanup(self) -> None:
        """
        清理输出目录 / Clean up output directory
        """
        if self.output_dir.exists():
            shutil.rmtree(self.output_dir)
            logger.info(f"输出目录已清理 / Output directory cleaned: {self.output_dir}")


# ==============================================================================
# 文档生成
# ==============================================================================

def _sanitize_filename(name: str, fallback: str = "document") -> str:
    base = re.sub(r"[^A-Za-z0-9._-]+", "_", str(name or ""))
    base = base.strip("._-")
    if not base:
        base = fallback
    return base[:64]

def _normalize_document_content(content: str) -> str:
    """
    轻量归一化文档内容，提升排版一致性
    """
    if not content:
        return ""
    text = content.replace("\r\n", "\n").replace("\r", "\n")
    lines = text.split("\n")
    output: List[str] = []
    in_code = False
    for raw in lines:
        stripped = raw.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            output.append(raw.rstrip())
            continue
        if in_code:
            output.append(raw.rstrip("\n"))
            continue
        indent_match = re.match(r"^(\s*)", raw)
        indent = indent_match.group(1) if indent_match else ""
        body = raw[len(indent):].rstrip()
        bullet_match = re.match(r"^[•·●▪▶◦]\s*(.+)$", body)
        if bullet_match:
            body = f"- {bullet_match.group(1).strip()}"
            output.append(indent + body)
            continue
        number_match = re.match(r"^(\d+)[、\)]\s*(.+)$", body)
        if number_match:
            body = f"{number_match.group(1)}. {number_match.group(2).strip()}"
            output.append(indent + body)
            continue
        output.append(indent + body)
    text = "\n".join(output)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text


def _split_table_row(line: str) -> List[str]:
    trimmed = line.strip().strip("|")
    if not trimmed:
        return []
    return [cell.strip() for cell in trimmed.split("|")]


def _is_table_separator(line: str) -> bool:
    trimmed = line.strip().strip("|")
    if not trimmed:
        return False
    parts = [part.strip() for part in trimmed.split("|")]
    for part in parts:
        if not re.match(r"^:?-{3,}:?$", part):
            return False
    return True


def _is_table_row(line: str) -> bool:
    return line.count("|") >= 2


def _parse_markdown_blocks(content: str) -> List[Dict[str, Any]]:
    lines = _normalize_document_content(content).split("\n")
    blocks: List[Dict[str, Any]] = []
    in_code = False
    code_lang = ""
    code_lines: List[str] = []
    quote_lines: List[str] = []
    paragraph_lines: List[str] = []
    table_lines: List[str] = []

    def flush_paragraph() -> None:
        if not paragraph_lines:
            return
        text = " ".join(paragraph_lines).strip()
        if text:
            blocks.append({"type": "para", "text": text})
        paragraph_lines.clear()

    def flush_quote() -> None:
        if not quote_lines:
            return
        text = "\n".join(quote_lines).strip()
        if text:
            blocks.append({"type": "quote", "text": text})
        quote_lines.clear()

    def flush_code() -> None:
        nonlocal in_code, code_lang
        if code_lines:
            blocks.append({"type": "code", "text": "\n".join(code_lines), "lang": code_lang})
        code_lines.clear()
        code_lang = ""
        in_code = False

    def flush_table() -> None:
        if not table_lines:
            return
        rows = [_split_table_row(row) for row in table_lines if not _is_table_separator(row)]
        header = None
        if len(table_lines) > 1 and _is_table_separator(table_lines[1]):
            header = rows[0] if rows else None
            rows = rows[1:] if len(rows) > 1 else []
        if rows or header:
            blocks.append({"type": "table", "header": header, "rows": rows})
        table_lines.clear()

    for raw in lines:
        stripped = raw.strip()
        if stripped.startswith("```"):
            flush_paragraph()
            flush_quote()
            flush_table()
            if in_code:
                flush_code()
            else:
                in_code = True
                code_lang = stripped[3:].strip()
            continue
        if in_code:
            code_lines.append(raw.rstrip("\n"))
            continue
        if _is_table_row(raw):
            flush_paragraph()
            flush_quote()
            table_lines.append(raw)
            continue
        if table_lines:
            flush_table()
        if stripped.startswith(">"):
            flush_paragraph()
            quote_lines.append(stripped.lstrip(">").strip())
            continue
        if quote_lines:
            flush_quote()
        if stripped and re.match(r"^(-{3,}|\*{3,}|_{3,})$", stripped):
            flush_paragraph()
            blocks.append({"type": "hr"})
            continue
        if not stripped:
            flush_paragraph()
            blocks.append({"type": "blank", "text": ""})
            continue
        if stripped.startswith("#"):
            flush_paragraph()
            level = len(stripped) - len(stripped.lstrip("#"))
            text = stripped[level:].strip() or "Untitled"
            blocks.append({"type": "heading", "level": min(level, 3), "text": text})
            continue
        bullet_match = re.match(r"^[-*]\s+(.*)$", stripped)
        if bullet_match:
            flush_paragraph()
            blocks.append({"type": "bullet", "text": bullet_match.group(1).strip(), "marker": "-"})
            continue
        number_match = re.match(r"^(\d+)\.\s+(.*)$", stripped)
        if number_match:
            flush_paragraph()
            blocks.append({
                "type": "number",
                "text": number_match.group(2).strip(),
                "marker": number_match.group(1),
            })
            continue
        paragraph_lines.append(stripped)

    if in_code:
        flush_code()
    if quote_lines:
        flush_quote()
    if table_lines:
        flush_table()
    if paragraph_lines:
        flush_paragraph()
    return blocks


def _split_reportlab_word(word: str, font_name: str, font_size: int, max_width: float) -> List[str]:
    if not word:
        return []
    current = ""
    parts = []
    for char in word:
        if pdfmetrics.stringWidth(current + char, font_name, font_size) <= max_width:
            current += char
        else:
            if current:
                parts.append(current)
            current = char
    if current:
        parts.append(current)
    return parts


def _wrap_reportlab_text(text: str, font_name: str, font_size: int, max_width: float) -> List[str]:
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = ""
    for word in words:
        chunks = _split_reportlab_word(word, font_name, font_size, max_width)
        for chunk in chunks:
            if not current:
                current = chunk
                continue
            trial = f"{current} {chunk}"
            if pdfmetrics.stringWidth(trial, font_name, font_size) <= max_width:
                current = trial
            else:
                lines.append(current)
                current = chunk
    if current:
        lines.append(current)
    return lines


def _wrap_reportlab_preformatted(text: str, font_name: str, font_size: int, max_width: float) -> List[str]:
    if text is None:
        return [""]
    lines: List[str] = []
    raw_lines = text.splitlines() or [""]
    for raw in raw_lines:
        current = ""
        for char in raw:
            if pdfmetrics.stringWidth(current + char, font_name, font_size) <= max_width:
                current += char
            else:
                lines.append(current)
                current = char
        lines.append(current)
    return lines


def _pick_reportlab_font(content: str) -> str:
    """
    尝试选择支持更多字符集的字体
    """
    if not REPORTLAB_AVAILABLE:
        return "Helvetica"
    text = content or ""
    if not any(ord(ch) > 127 for ch in text):
        return "Helvetica"
    font_dir = os.environ.get("WINDIR", "C:\\Windows")
    font_dir = os.path.join(font_dir, "Fonts")
    candidates = [
        ("SamaHei", os.path.join(font_dir, "simhei.ttf")),
        ("SamaFang", os.path.join(font_dir, "simfang.ttf")),
        ("SamaKai", os.path.join(font_dir, "simkai.ttf")),
    ]
    registered = set(pdfmetrics.getRegisteredFontNames())
    for name, path in candidates:
        if not os.path.exists(path):
            continue
        if name in registered:
            return name
        try:
            pdfmetrics.registerFont(TTFont(name, path))
            return name
        except Exception:
            continue
    return "Helvetica"


def _generate_pdf_reportlab(content: str, output_path: Path, title: Optional[str]) -> None:
    canvas_obj = canvas.Canvas(str(output_path), pagesize=A4)
    page_width, page_height = A4
    margin = 54
    base_size = 11
    line_gap = 4
    paragraph_gap = 6
    font_name = _pick_reportlab_font(content)
    code_font = "Courier"
    y = page_height - margin

    def new_page() -> None:
        nonlocal y
        canvas_obj.showPage()
        canvas_obj.setFont(font_name, base_size)
        y = page_height - margin

    def ensure_space(height: float) -> None:
        nonlocal y
        if y - height < margin:
            new_page()

    blocks = _parse_markdown_blocks(content)
    if title:
        blocks.insert(0, {"type": "heading", "level": 1, "text": title})
        blocks.insert(1, {"type": "blank", "text": ""})

    for block in blocks:
        block_type = block.get("type")
        if block_type == "blank":
            y -= base_size + line_gap
            if y < margin:
                new_page()
            continue
        if block_type == "heading":
            level = int(block.get("level") or 1)
            size = 16 - (level - 1) * 2
            size = max(size, 12)
            canvas_obj.setFont(font_name, size)
            lines = _wrap_reportlab_text(block.get("text", ""), font_name, size, page_width - margin * 2)
            for line in lines:
                if y < margin + size:
                    new_page()
                canvas_obj.drawString(margin, y, line)
                y -= size + line_gap
            y -= line_gap
            canvas_obj.setFont(font_name, base_size)
            continue
        if block_type == "hr":
            if y < margin + base_size:
                new_page()
            canvas_obj.setLineWidth(1)
            canvas_obj.line(margin, y, page_width - margin, y)
            y -= base_size + line_gap
            continue
        if block_type == "quote":
            quote_indent = 16
            lines = _wrap_reportlab_text(block.get("text", ""), font_name, base_size, page_width - margin * 2 - quote_indent)
            canvas_obj.setFont(font_name, base_size)
            y_start = y
            for line in lines:
                if y < margin + base_size:
                    new_page()
                    canvas_obj.setFont(font_name, base_size)
                canvas_obj.drawString(margin + quote_indent, y, line)
                y -= base_size + line_gap
            y_end = y + line_gap
            canvas_obj.setStrokeColor(colors.grey)
            canvas_obj.setLineWidth(1)
            canvas_obj.line(margin + 4, y_start, margin + 4, y_end)
            canvas_obj.setStrokeColor(colors.black)
            canvas_obj.setFont(font_name, base_size)
            y -= paragraph_gap
            continue
        if block_type == "code":
            code_indent = 12
            lines = _wrap_reportlab_preformatted(
                block.get("text", ""),
                code_font,
                base_size,
                page_width - margin * 2 - code_indent,
            )
            line_height = base_size + line_gap
            block_height = max(1, len(lines)) * line_height + 8
            ensure_space(block_height)
            canvas_obj.setFillColor(colors.whitesmoke)
            canvas_obj.rect(margin, y - block_height, page_width - margin * 2, block_height, fill=1, stroke=0)
            canvas_obj.setFillColor(colors.black)
            canvas_obj.setFont(code_font, base_size)
            text_y = y - 6 - base_size
            for line in lines:
                canvas_obj.drawString(margin + code_indent, text_y, line)
                text_y -= line_height
            y -= block_height
            canvas_obj.setFont(font_name, base_size)
            y -= paragraph_gap
            continue
        if block_type == "table":
            header = block.get("header") or []
            rows = block.get("rows") or []
            table_rows = [header] + rows if header else rows
            if not table_rows:
                continue
            col_count = max(len(row) for row in table_rows) if table_rows else 0
            if col_count <= 0:
                continue
            col_width = (page_width - margin * 2) / col_count
            line_height = base_size + line_gap
            for row_index, row in enumerate(table_rows):
                cells = list(row) + [""] * (col_count - len(row))
                wrapped = [
                    _wrap_reportlab_text(cell, font_name, base_size, col_width - 8)
                    for cell in cells
                ]
                row_height = max(len(lines) for lines in wrapped) * line_height + 8
                ensure_space(row_height)
                if header and row_index == 0:
                    canvas_obj.setFillColor(colors.lightgrey)
                    canvas_obj.rect(margin, y - row_height, page_width - margin * 2, row_height, fill=1, stroke=0)
                    canvas_obj.setFillColor(colors.black)
                for col_index, lines in enumerate(wrapped):
                    cell_x = margin + col_index * col_width
                    canvas_obj.rect(cell_x, y - row_height, col_width, row_height, stroke=1, fill=0)
                    text_y = y - 6 - base_size
                    for line in lines:
                        canvas_obj.drawString(cell_x + 4, text_y, line)
                        text_y -= line_height
                y -= row_height
            y -= paragraph_gap
            continue
        if block_type in {"bullet", "number"}:
            marker = block.get("marker") or "-"
            prefix = f"{marker} "
            text = prefix + block.get("text", "")
            indent = 10
        else:
            text = block.get("text", "")
            indent = 0
        canvas_obj.setFont(font_name, base_size)
        lines = _wrap_reportlab_text(text, font_name, base_size, page_width - margin * 2 - indent)
        for line in lines:
            if y < margin + base_size:
                new_page()
            canvas_obj.drawString(margin + indent, y, line)
            y -= base_size + line_gap
        y -= paragraph_gap

    canvas_obj.save()


def _split_fitz_word(word: str, font: "fitz.Font", font_size: int, max_width: float) -> List[str]:
    if not word:
        return []
    current = ""
    parts = []
    for char in word:
        if font.text_length(current + char, font_size) <= max_width:
            current += char
        else:
            if current:
                parts.append(current)
            current = char
    if current:
        parts.append(current)
    return parts


def _wrap_fitz_text(text: str, font: "fitz.Font", font_size: int, max_width: float) -> List[str]:
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = ""
    for word in words:
        chunks = _split_fitz_word(word, font, font_size, max_width)
        for chunk in chunks:
            if not current:
                current = chunk
                continue
            trial = f"{current} {chunk}"
            if font.text_length(trial, font_size) <= max_width:
                current = trial
            else:
                lines.append(current)
                current = chunk
    if current:
        lines.append(current)
    return lines


def _wrap_fitz_preformatted(text: str, font: "fitz.Font", font_size: int, max_width: float) -> List[str]:
    if text is None:
        return [""]
    lines: List[str] = []
    raw_lines = text.splitlines() or [""]
    for raw in raw_lines:
        current = ""
        for char in raw:
            if font.text_length(current + char, font_size) <= max_width:
                current += char
            else:
                lines.append(current)
                current = char
        lines.append(current)
    return lines


def _pick_fitz_font(content: str) -> "fitz.Font":
    """
    尝试为 PyMuPDF 选择支持更多字符集的字体
    """
    if not FITZ_AVAILABLE:
        return fitz.Font("helv")
    text = content or ""
    if not any(ord(ch) > 127 for ch in text):
        return fitz.Font("helv")
    font_dir = os.environ.get("WINDIR", "C:\\Windows")
    font_dir = os.path.join(font_dir, "Fonts")
    candidates = [
        os.path.join(font_dir, "simhei.ttf"),
        os.path.join(font_dir, "simsun.ttc"),
        os.path.join(font_dir, "msyh.ttc"),
    ]
    for path in candidates:
        if not os.path.exists(path):
            continue
        try:
            return fitz.Font(fontfile=path)
        except Exception:
            continue
    return fitz.Font("helv")


def _generate_pdf_fitz(content: str, output_path: Path, title: Optional[str]) -> None:
    page_width = 595
    page_height = 842
    margin = 54
    base_size = 11
    line_gap = 4
    paragraph_gap = 6
    font = _pick_fitz_font(content)
    mono_font = fitz.Font("cour")
    doc = fitz.open()
    page = doc.new_page(width=page_width, height=page_height)
    y = margin

    def new_page() -> None:
        nonlocal page, y
        page = doc.new_page(width=page_width, height=page_height)
        y = margin

    blocks = _parse_markdown_blocks(content)
    if title:
        blocks.insert(0, {"type": "heading", "level": 1, "text": title})
        blocks.insert(1, {"type": "blank", "text": ""})

    for block in blocks:
        block_type = block.get("type")
        if block_type == "blank":
            y += base_size + line_gap
            if y > page_height - margin:
                new_page()
            continue
        if block_type == "heading":
            level = int(block.get("level") or 1)
            size = 16 - (level - 1) * 2
            size = max(size, 12)
            lines = _wrap_fitz_text(block.get("text", ""), font, size, page_width - margin * 2)
            for line in lines:
                if y > page_height - margin:
                    new_page()
                page.insert_text((margin, y), line, fontsize=size, fontname=font.name)
                y += size + line_gap
            y += line_gap
            continue
        if block_type == "hr":
            if y > page_height - margin:
                new_page()
            page.draw_line((margin, y), (page_width - margin, y), width=1)
            y += base_size + line_gap
            continue
        if block_type == "quote":
            quote_indent = 16
            lines = _wrap_fitz_text(block.get("text", ""), font, base_size, page_width - margin * 2 - quote_indent)
            block_height = max(1, len(lines)) * (base_size + line_gap)
            if y + block_height > page_height - margin:
                new_page()
            page.draw_line((margin + 4, y), (margin + 4, y + block_height - line_gap), color=(0.5, 0.5, 0.5), width=1)
            for line in lines:
                if y > page_height - margin:
                    new_page()
                page.insert_text((margin + quote_indent, y), line, fontsize=base_size, fontname=font.name)
                y += base_size + line_gap
            y += paragraph_gap
            continue
        if block_type == "code":
            code_indent = 12
            lines = _wrap_fitz_preformatted(block.get("text", ""), mono_font, base_size, page_width - margin * 2 - code_indent)
            block_height = max(1, len(lines)) * (base_size + line_gap) + 8
            if y + block_height > page_height - margin:
                new_page()
            rect = fitz.Rect(margin, y - 4, page_width - margin, y + block_height - 4)
            page.draw_rect(rect, color=None, fill=(0.95, 0.95, 0.95))
            for line in lines:
                if y > page_height - margin:
                    new_page()
                page.insert_text((margin + code_indent, y), line, fontsize=base_size, fontname="cour")
                y += base_size + line_gap
            y += paragraph_gap
            continue
        if block_type == "table":
            header = block.get("header") or []
            rows = block.get("rows") or []
            table_rows = [header] + rows if header else rows
            if not table_rows:
                continue
            for row_index, row in enumerate(table_rows):
                line = " | ".join([cell for cell in row if cell is not None])
                lines = _wrap_fitz_text(line, font, base_size, page_width - margin * 2)
                for text in lines:
                    if y > page_height - margin:
                        new_page()
                    page.insert_text((margin, y), text, fontsize=base_size, fontname=font.name)
                    y += base_size + line_gap
                if header and row_index == 0:
                    page.draw_line((margin, y), (page_width - margin, y), width=1)
                    y += line_gap
            y += paragraph_gap
            continue
        if block_type in {"bullet", "number"}:
            marker = block.get("marker") or "-"
            prefix = f"{marker} "
            text = prefix + block.get("text", "")
            indent = 10
        else:
            text = block.get("text", "")
            indent = 0
        lines = _wrap_fitz_text(text, font, base_size, page_width - margin * 2 - indent)
        for line in lines:
            if y > page_height - margin:
                new_page()
            page.insert_text((margin + indent, y), line, fontsize=base_size, fontname=font.name)
            y += base_size + line_gap
        y += paragraph_gap

    doc.save(str(output_path))
    doc.close()


def generate_docx(content: str, output_path: Path, title: Optional[str] = None) -> None:
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    from docx.shared import Inches, Pt, RGBColor
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    document = docx.Document()
    section = document.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)
    normal_style.paragraph_format.line_spacing = 1.2
    normal_style.paragraph_format.space_after = Pt(6)

    def set_run_font(run, font_name: str, font_size: int, bold: bool = False, italic: bool = False) -> None:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        rpr = run._element.rPr
        if rpr is None:
            rpr = OxmlElement("w:rPr")
            run._element.append(rpr)
        r_fonts = rpr.find(qn("w:rFonts"))
        if r_fonts is None:
            r_fonts = OxmlElement("w:rFonts")
            rpr.append(r_fonts)
        r_fonts.set(qn("w:ascii"), font_name)
        r_fonts.set(qn("w:hAnsi"), font_name)
        r_fonts.set(qn("w:eastAsia"), font_name)

    def set_paragraph_shading(paragraph, fill: str = "F2F2F2") -> None:
        p = paragraph._p
        p_pr = p.get_or_add_pPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), fill)
        p_pr.append(shd)

    def set_paragraph_border(paragraph, color: str = "D0D0D0") -> None:
        p = paragraph._p
        p_pr = p.get_or_add_pPr()
        p_bdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "4")
        bottom.set(qn("w:color"), color)
        p_bdr.append(bottom)
        p_pr.append(p_bdr)

    if title:
        document.core_properties.title = title
        heading = document.add_heading(title, level=1)
        heading.paragraph_format.space_after = Pt(10)
    for block in _parse_markdown_blocks(content):
        block_type = block.get("type")
        if block_type == "heading":
            paragraph = document.add_heading(block.get("text", ""), level=int(block.get("level") or 1))
            paragraph.paragraph_format.space_before = Pt(12)
            paragraph.paragraph_format.space_after = Pt(6)
        elif block_type == "bullet":
            paragraph = document.add_paragraph(block.get("text", ""), style="List Bullet")
            paragraph.paragraph_format.space_after = Pt(2)
        elif block_type == "number":
            paragraph = document.add_paragraph(block.get("text", ""), style="List Number")
            paragraph.paragraph_format.space_after = Pt(2)
        elif block_type == "quote":
            paragraph = document.add_paragraph(block.get("text", ""))
            try:
                paragraph.style = "Intense Quote"
            except (KeyError, ValueError):
                paragraph.paragraph_format.left_indent = docx.shared.Pt(18)
                paragraph.paragraph_format.space_after = docx.shared.Pt(6)
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(90, 90, 90)
                run.font.italic = True
        elif block_type == "code":
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.left_indent = Pt(12)
            paragraph.paragraph_format.space_after = docx.shared.Pt(6)
            paragraph.paragraph_format.space_before = docx.shared.Pt(4)
            set_paragraph_shading(paragraph)
            lines = (block.get("text", "") or "").splitlines() or [""]
            for idx, line in enumerate(lines):
                run = paragraph.add_run(line)
                set_run_font(run, "Consolas", 10)
                if idx < len(lines) - 1:
                    run.add_break()
        elif block_type == "hr":
            paragraph = document.add_paragraph("")
            set_paragraph_border(paragraph)
        elif block_type == "table":
            header = block.get("header") or []
            rows = block.get("rows") or []
            table_rows = [header] + rows if header else rows
            if not table_rows:
                continue
            col_count = max(len(row) for row in table_rows)
            table = document.add_table(rows=1, cols=col_count)
            table.style = "Table Grid"
            table.autofit = True
            first_row = table.rows[0].cells
            header_row = header if header else table_rows[0]
            for idx, text in enumerate(list(header_row) + [""] * (col_count - len(header_row))):
                run = first_row[idx].paragraphs[0].add_run(text)
                run.bold = True if header else False
            data_rows = rows if header else table_rows[1:]
            for row in data_rows:
                cells = list(row) + [""] * (col_count - len(row))
                row_cells = table.add_row().cells
                for cell_idx, text in enumerate(cells):
                    row_cells[cell_idx].text = str(text)
        elif block_type == "blank":
            document.add_paragraph("")
        else:
            document.add_paragraph(block.get("text", ""))
    document.save(str(output_path))


def generate_pdf(content: str, output_path: Path, title: Optional[str] = None) -> None:
    if REPORTLAB_AVAILABLE:
        _generate_pdf_reportlab(content, output_path, title)
        return
    if FITZ_AVAILABLE:
        _generate_pdf_fitz(content, output_path, title)
        return
    raise ImportError("reportlab or PyMuPDF not installed. Run: pip install reportlab or pip install PyMuPDF")


def generate_document(
    content: str,
    doc_type: str,
    output_dir: str,
    title: Optional[str] = None,
    filename: Optional[str] = None
) -> Dict[str, Any]:
    if not content:
        raise ValueError("Content required.")
    doc_type = str(doc_type or "").lower()
    if doc_type not in {"pdf", "docx"}:
        raise ValueError("Unsupported format.")
    safe_title = title.strip() if isinstance(title, str) else ""
    base_name = _sanitize_filename(filename or safe_title or "document")
    stamp = time.strftime("%Y%m%d_%H%M%S")
    task_id = f"doc_{stamp}_{uuid.uuid4().hex[:6]}"
    output_path = Path(output_dir) / task_id
    output_path.mkdir(parents=True, exist_ok=True)
    file_name = f"{base_name}.{doc_type}"
    file_path = output_path / file_name
    if doc_type == "docx":
        generate_docx(content, file_path, safe_title or None)
    else:
        generate_pdf(content, file_path, safe_title or None)
    return {
        "task_id": task_id,
        "file_name": file_name,
        "file_path": str(file_path),
        "title": safe_title or base_name,
        "format": doc_type,
    }


# ==============================================================================
# 便捷函数
# ==============================================================================

def preprocess_files(
    task_id: str,
    file_paths: List[str],
    output_dir: Optional[str] = None
) -> Dict[str, Any]:
    """
    预处理文件列表 / Preprocess file list
    
    便捷函数，创建DocumentConverter并处理文件
    Convenience function that creates DocumentConverter and processes files
    
    Args:
        task_id: 任务ID / Task ID
        file_paths: 文件路径列表 / File path list
        output_dir: 输出目录 / Output directory
        
    Returns:
        Dict[str, Any]: 处理结果 / Processing result
    """
    converter = DocumentConverter(task_id=task_id, output_dir=output_dir)
    return converter.process(file_paths)


def get_supported_extensions() -> Dict[str, List[str]]:
    """
    获取支持的文件扩展名 / Get supported file extensions
    
    Returns:
        Dict[str, List[str]]: 按类型分组的扩展名 / Extensions grouped by type
    """
    return DocumentConverter.SUPPORTED_EXTENSIONS.copy()


def is_file_supported(file_path: str) -> bool:
    """
    检查文件是否受支持 / Check if file is supported
    
    Args:
        file_path: 文件路径 / File path
        
    Returns:
        bool: 是否支持 / Whether supported
    """
    ext = Path(file_path).suffix.lower()
    all_exts = []
    for exts in DocumentConverter.SUPPORTED_EXTENSIONS.values():
        all_exts.extend(exts)
    return ext in all_exts
